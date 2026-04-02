"""
Parses an uploaded RFP (PDF or Word) and extracts structured information.
"""
import os
import tempfile
from pathlib import Path
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_rfp_text(uploaded_file) -> str:
    """
    Extract raw text from an uploaded RFP file.
    Works with PDF, DOCX, PPTX.
    uploaded_file: Streamlit UploadedFile object
    """
    suffix = Path(uploaded_file.name).suffix.lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.read())
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            return _extract_pdf(tmp_path)
        elif suffix in (".docx", ".doc"):
            return _extract_docx(tmp_path)
        elif suffix in (".pptx", ".ppt"):
            return _extract_pptx(tmp_path)
        else:
            return ""
    finally:
        os.unlink(tmp_path)


def _extract_pdf(path: str) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(pages)


def _extract_docx(path: str) -> str:
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        content.append(para.text.strip())
        if content:
            slides.append(f"[Slide {i}]\n" + "\n".join(content))
    return "\n\n".join(slides)
