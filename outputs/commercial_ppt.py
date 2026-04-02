"""
Generates the Commercial Proposal PowerPoint.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from config import OUTPUTS_PATH, BLENDED_RATE_USD
from outputs.ppt_generator import (
    _blank_slide, _set_bg, _add_textbox, _add_slide_header,
    _add_bullet_column, _add_closing_slide,
    RED, DARK_GRAY, MID_GRAY, LIGHT_GRAY, WHITE,
    SLIDE_W, SLIDE_H
)

PHASE_COLORS_RGB = [
    RGBColor(0xD6, 0xE4, 0xF0),
    RGBColor(0xD5, 0xF5, 0xE3),
    RGBColor(0xFE, 0xF9, 0xE7),
    RGBColor(0xFD, 0xED, 0xEC),
]


def generate_commercial_ppt(proposal_data: dict, client_name: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    commercial = proposal_data.get("commercial", {})
    effort = proposal_data.get("effort_model", {})
    rfp = proposal_data.get("rfp_intel", {})

    _add_cover_commercial(prs, commercial.get("cover", {}), client_name)
    _add_commercial_summary(prs, commercial.get("commercial_summary", {}), effort)
    _add_fee_by_phase(prs, commercial.get("fee_by_phase", []), effort)
    _add_fee_detail(prs, effort)
    _add_payment_milestones(prs, commercial.get("payment_milestones", []))
    _add_assumptions_exclusions(prs, commercial)
    _add_terms_slide(prs, commercial.get("terms", {}))
    _add_closing_slide(prs, client_name)

    filename = f"Commercial_Proposal_{client_name.replace(' ', '_')}.pptx"
    output_path = os.path.join(OUTPUTS_PATH, filename)
    prs.save(output_path)
    return output_path


def _add_cover_commercial(prs, cover: dict, client_name: str):
    slide = _blank_slide(prs)
    _set_bg(slide, DARK_GRAY)

    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RED
    left_bar.line.fill.background()

    _add_textbox(slide, cover.get("title", "Commercial Proposal"),
                 Inches(0.8), Inches(1.5), Inches(11), Inches(2),
                 font_size=36, bold=True, color=WHITE)
    _add_textbox(slide, "Strictly Confidential",
                 Inches(0.8), Inches(3.8), Inches(9), Inches(0.5),
                 font_size=14, color=RGBColor(0xFF, 0x80, 0x80))
    _add_textbox(slide, f"Prepared for: {client_name}",
                 Inches(0.8), Inches(4.4), Inches(9), Inches(0.5),
                 font_size=14, color=WHITE)
    _add_textbox(slide, cover.get("date", ""),
                 Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
                 font_size=12, color=MID_GRAY)
    _add_textbox(slide, "PROTIVITI MIDDLE EAST | REAL ESTATE & INFRASTRUCTURE PRACTICE",
                 Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
                 font_size=9, color=MID_GRAY)


def _add_commercial_summary(prs, summary: dict, effort: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Commercial Summary")

    total_fee = effort.get("total_fee_usd", summary.get("total_fee_usd", 0))
    total_hours = effort.get("total_hours", summary.get("total_hours", 0))
    duration = summary.get("duration_weeks", "TBD")
    payment = summary.get("payment_structure", "Milestone-based")
    validity = summary.get("validity_days", 30)

    # Key metrics boxes
    metrics = [
        ("Total Fee", f"USD {total_fee:,.0f}"),
        ("Total Effort", f"{total_hours:,} hours"),
        ("Duration", f"{duration} weeks"),
        ("Payment Structure", payment),
        ("Proposal Validity", f"{validity} days"),
    ]
    box_w = Inches(2.3)
    for i, (label, value) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(2.5)
        box = slide.shapes.add_shape(1, x, Inches(1.8), box_w, Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RED
        _add_textbox(slide, label, x + Inches(0.1), Inches(1.9),
                     box_w - Inches(0.2), Inches(0.4),
                     font_size=9, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, value, x + Inches(0.1), Inches(2.4),
                     box_w - Inches(0.2), Inches(1.0),
                     font_size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

    _add_textbox(slide, f"Blended Rate: USD {BLENDED_RATE_USD}/hour",
                 Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
                 font_size=10, color=MID_GRAY)


def _add_fee_by_phase(prs, fee_phases: list, effort: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Fee by Phase")

    phases_data = effort.get("phases", fee_phases)
    total_fee = effort.get("total_fee_usd", 1)

    y = Inches(1.6)
    for i, phase in enumerate(phases_data):
        phase_name = phase.get("phase_name", phase.get("phase", ""))
        phase_fee = phase.get("phase_fee_usd", phase.get("fee_usd", 0))
        phase_hours = phase.get("phase_hours", phase.get("hours", 0))
        pct = (phase_fee / total_fee * 100) if total_fee else 0

        color = PHASE_COLORS_RGB[i % len(PHASE_COLORS_RGB)]
        box = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        _add_textbox(slide, phase_name, Inches(0.6), y + Inches(0.1),
                     Inches(6), Inches(0.5), font_size=11, bold=True, color=DARK_GRAY)
        _add_textbox(slide, f"{phase_hours:,} hrs", Inches(7), y + Inches(0.1),
                     Inches(2), Inches(0.5), font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, f"USD {phase_fee:,.0f}", Inches(9.2), y + Inches(0.1),
                     Inches(2), Inches(0.5), font_size=11, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, f"{pct:.0f}%", Inches(11.4), y + Inches(0.1),
                     Inches(1.2), Inches(0.5), font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        y += Inches(0.85)

    # Total row
    y += Inches(0.2)
    _add_textbox(slide, f"TOTAL  |  USD {total_fee:,.0f}",
                 Inches(0.5), y, Inches(12.3), Inches(0.5),
                 font_size=14, bold=True, color=RED, align=PP_ALIGN.RIGHT)


def _add_fee_detail(prs, effort: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Fee by Deliverable (L2 Detail)")

    _add_textbox(slide, "See Annexure A: Detailed Effort & Cost Model (Excel) for full breakdown",
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
                 font_size=11, color=MID_GRAY)

    # Summary table of L1 deliverables
    headers = ["Phase", "L1 Deliverable", "Hours", "Fee (USD)"]
    col_x = [Inches(0.5), Inches(3.0), Inches(9.5), Inches(10.8)]
    col_w = [Inches(2.4), Inches(6.4), Inches(1.2), Inches(1.8)]

    y = Inches(2.2)
    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        _add_textbox(slide, h, x, y, w, Inches(0.35),
                     font_size=9, bold=True, color=WHITE)
        box = slide.shapes.add_shape(1, x, y, w, Inches(0.35))
        box.fill.solid()
        box.fill.fore_color.rgb = RED
        box.line.fill.background()

    y += Inches(0.4)
    for p_idx, phase in enumerate(effort.get("phases", [])):
        for deliv in phase.get("deliverables", []):
            _add_textbox(slide, phase.get("phase_name", ""), col_x[0], y,
                         col_w[0], Inches(0.3), font_size=8.5, color=DARK_GRAY)
            _add_textbox(slide, deliv.get("l1_name", ""), col_x[1], y,
                         col_w[1], Inches(0.3), font_size=8.5, color=DARK_GRAY)
            _add_textbox(slide, str(deliv.get("l1_hours", "")), col_x[2], y,
                         col_w[2], Inches(0.3), font_size=8.5, color=DARK_GRAY, align=PP_ALIGN.CENTER)
            _add_textbox(slide, f"USD {deliv.get('l1_fee_usd', 0):,.0f}", col_x[3], y,
                         col_w[3], Inches(0.3), font_size=8.5, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
            y += Inches(0.32)


def _add_payment_milestones(prs, milestones: list):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Payment Milestones")

    y = Inches(1.6)
    for i, m in enumerate(milestones):
        color = PHASE_COLORS_RGB[i % len(PHASE_COLORS_RGB)]
        box = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        _add_textbox(slide, f"#{i+1}", Inches(0.6), y + Inches(0.1),
                     Inches(0.5), Inches(0.6), font_size=14, bold=True, color=RED)
        _add_textbox(slide, m.get("milestone", ""), Inches(1.2), y + Inches(0.05),
                     Inches(5), Inches(0.4), font_size=11, bold=True, color=DARK_GRAY)
        _add_textbox(slide, m.get("trigger", ""), Inches(1.2), y + Inches(0.45),
                     Inches(5.5), Inches(0.35), font_size=9.5, color=MID_GRAY)
        _add_textbox(slide, f"Week {m.get('due_week', '')}", Inches(7.2), y + Inches(0.2),
                     Inches(2), Inches(0.4), font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        amount = m.get("amount_usd", 0)
        _add_textbox(slide, f"USD {amount:,.0f}", Inches(9.5), y + Inches(0.1),
                     Inches(2.0), Inches(0.5), font_size=13, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, m.get("percentage", ""), Inches(11.7), y + Inches(0.2),
                     Inches(0.9), Inches(0.4), font_size=10, color=DARK_GRAY)
        y += Inches(1.0)


def _add_assumptions_exclusions(prs, commercial: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Key Assumptions & Exclusions")

    _add_bullet_column(slide, "Key Assumptions",
                       commercial.get("key_assumptions", []),
                       Inches(0.5), Inches(1.6), Inches(5.8))
    _add_bullet_column(slide, "Exclusions",
                       commercial.get("exclusions", []),
                       Inches(6.8), Inches(1.6), Inches(5.8))


def _add_terms_slide(prs, terms: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, "Terms of Engagement")

    items = [
        ("Payment Terms", terms.get("payment_terms", "")),
        ("Variation Process", terms.get("variation_process", "")),
        ("IP Ownership", terms.get("ip_ownership", "")),
        ("Confidentiality", terms.get("confidentiality", "")),
    ]
    y = Inches(1.6)
    for label, value in items:
        _add_textbox(slide, label, Inches(0.5), y, Inches(3.5), Inches(0.35),
                     font_size=11, bold=True, color=RED)
        _add_textbox(slide, value, Inches(4.0), y, Inches(9), Inches(0.6),
                     font_size=10.5, color=DARK_GRAY)
        y += Inches(1.2)
