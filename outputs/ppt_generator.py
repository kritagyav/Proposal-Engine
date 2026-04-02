"""
Generates the Technical Proposal PowerPoint.
Uses python-pptx to build a professionally structured deck.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from config import OUTPUTS_PATH, TEMPLATES_PATH

# Protiviti brand colors
RED = RGBColor(0xC8, 0x10, 0x2E)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
MID_GRAY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def generate_technical_ppt(proposal_data: dict, client_name: str) -> str:
    """
    Generate the technical proposal PPT.
    Returns the file path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    tech = proposal_data.get("technical", {})
    rfp = proposal_data.get("rfp_intel", {})
    web_research = proposal_data.get("web_research", {})
    value_add_slide = proposal_data.get("value_add_slide", {})

    _add_cover_slide(prs, tech.get("cover", {}))
    _add_section_divider(prs, "01", "Executive Summary")
    _add_executive_summary(prs, tech.get("executive_summary", {}))

    # Client overview — powered by web research
    _add_section_divider(prs, "02", "Client Overview")
    _add_client_overview_slide(prs, web_research.get("client_profile", {}), rfp)
    _add_sector_context_slide(prs, web_research.get("sector_context", {}), rfp)

    _add_section_divider(prs, "03", "Our Understanding")
    _add_understanding_slide(prs, tech.get("our_understanding", {}))
    _add_section_divider(prs, "04", "Value Proposition")
    _add_value_proposition(prs, tech.get("value_proposition", {}))
    _add_past_relationship(prs, tech.get("past_relationship", {}))
    _add_section_divider(prs, "05", "Scope of Work")
    _add_scope_slides(prs, tech.get("scope_of_work", {}))

    # Value-add / leading practice recommendations
    if value_add_slide:
        _add_section_divider(prs, "06", "Our Recommended Enhancements")
        _add_value_add_slide(prs, value_add_slide)

    _add_section_divider(prs, "07", "Our Approach")
    _add_approach_slide(prs, tech.get("approach_methodology", {}))
    _add_section_divider(prs, "08", "Engagement Governance")
    _add_governance_slide(prs, tech.get("engagement_governance", {}))
    _add_section_divider(prs, "09", "Our Team")
    _add_team_slide(prs, tech.get("project_team", {}))
    _add_section_divider(prs, "10", "Timeline")
    _add_timeline_slide(prs, tech.get("timeline", {}))
    _add_section_divider(prs, "11", "Experience & Why Protiviti")
    _add_experience_slide(prs, tech.get("relevant_experience", {}))
    _add_why_protiviti(prs, tech.get("why_protiviti", {}))
    _add_assumptions_slide(prs, tech.get("key_assumptions", {}))
    _add_closing_slide(prs, client_name)

    filename = f"Technical_Proposal_{client_name.replace(' ', '_')}.pptx"
    output_path = os.path.join(OUTPUTS_PATH, filename)
    prs.save(output_path)
    return output_path


def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_cover_slide(prs, cover: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, DARK_GRAY)

    # Red accent bar on left
    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), SLIDE_H)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RED
    left_bar.line.fill.background()

    # Title
    _add_textbox(slide, cover.get("title", "Technical Proposal"),
                 Inches(0.8), Inches(1.5), Inches(11), Inches(2),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    _add_textbox(slide, cover.get("subtitle", "Technical Proposal | Confidential"),
                 Inches(0.8), Inches(3.8), Inches(9), Inches(0.6),
                 font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Client
    _add_textbox(slide, f"Prepared for: {cover.get('client', '')}",
                 Inches(0.8), Inches(4.6), Inches(9), Inches(0.5),
                 font_size=14, color=WHITE)

    # Date
    _add_textbox(slide, cover.get("date", ""),
                 Inches(0.8), Inches(5.2), Inches(6), Inches(0.4),
                 font_size=12, color=MID_GRAY)

    # Protiviti footer
    _add_textbox(slide, "PROTIVITI MIDDLE EAST | REAL ESTATE & INFRASTRUCTURE PRACTICE",
                 Inches(0.8), Inches(6.7), Inches(12), Inches(0.5),
                 font_size=9, color=MID_GRAY)


def _add_section_divider(prs, number: str, title: str):
    slide = _blank_slide(prs)
    _set_bg(slide, RED)
    _add_textbox(slide, number, Inches(1), Inches(2.5), Inches(2), Inches(1.5),
                 font_size=72, bold=True, color=RGBColor(0xFF, 0x80, 0x80))
    _add_textbox(slide, title, Inches(1), Inches(4.0), Inches(10), Inches(1.2),
                 font_size=36, bold=True, color=WHITE)


def _add_executive_summary(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Executive Summary"))

    content = [
        ("Our Understanding", data.get("our_understanding", "")),
        ("Our Approach", data.get("our_approach", "")),
        ("Our Commitment", data.get("our_commitment", "")),
    ]
    y = Inches(1.6)
    for label, text in content:
        _add_textbox(slide, label, Inches(0.5), y, Inches(3.5), Inches(0.35),
                     font_size=11, bold=True, color=RED)
        _add_textbox(slide, text, Inches(0.5), y + Inches(0.35), Inches(12.3), Inches(0.8),
                     font_size=10.5, color=DARK_GRAY)
        y += Inches(1.5)


def _add_understanding_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Understanding"))

    _add_bullet_column(slide, "Key Drivers", data.get("key_drivers", []),
                       Inches(0.5), Inches(1.6), Inches(4))
    _add_bullet_column(slide, "Challenges Identified", data.get("challenges_identified", []),
                       Inches(4.8), Inches(1.6), Inches(4))
    _add_bullet_column(slide, "Success Factors", data.get("success_factors", []),
                       Inches(9.0), Inches(1.6), Inches(4))


def _add_value_proposition(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Value Proposition"))

    # Headline
    _add_textbox(slide, data.get("headline", ""), Inches(0.5), Inches(1.5),
                 Inches(12.3), Inches(0.7), font_size=14, bold=True, color=RED)

    y = Inches(2.4)
    for vp in data.get("value_points", []):
        _add_textbox(slide, f"► {vp.get('point', '')}", Inches(0.5), y,
                     Inches(12.3), Inches(0.3), font_size=11, bold=True, color=DARK_GRAY)
        _add_textbox(slide, vp.get("detail", ""), Inches(0.8), y + Inches(0.3),
                     Inches(12), Inches(0.5), font_size=10, color=MID_GRAY)
        y += Inches(1.0)


def _add_past_relationship(prs, data: dict):
    if not data or not data.get("relationship_narrative"):
        return
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Relationship"))

    _add_textbox(slide, data.get("relationship_narrative", ""),
                 Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2),
                 font_size=11, color=DARK_GRAY)

    _add_bullet_column(slide, "Past Engagements", data.get("past_engagements", []),
                       Inches(0.5), Inches(3.2), Inches(6))

    _add_textbox(slide, "Continuity Benefit", Inches(7), Inches(3.0),
                 Inches(5.8), Inches(0.4), font_size=11, bold=True, color=RED)
    _add_textbox(slide, data.get("continuity_benefit", ""), Inches(7), Inches(3.5),
                 Inches(5.8), Inches(1.5), font_size=10.5, color=DARK_GRAY)


def _add_scope_slides(prs, scope: dict):
    phases = scope.get("phases", [])
    if not phases:
        return
    for phase in phases:
        slide = _blank_slide(prs)
        _set_bg(slide, WHITE)
        _add_slide_header(slide, phase.get("phase_name", "Scope of Work"))

        _add_textbox(slide, phase.get("phase_objective", ""),
                     Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                     font_size=11, color=MID_GRAY)

        y = Inches(2.2)
        for deliv in phase.get("deliverables_l1", []):
            _add_textbox(slide, f"■ {deliv.get('name', '')}", Inches(0.5), y,
                         Inches(12.3), Inches(0.35), font_size=11, bold=True, color=RED)
            _add_textbox(slide, deliv.get("description", ""), Inches(0.8), y + Inches(0.35),
                         Inches(12), Inches(0.3), font_size=10, color=DARK_GRAY)
            for sub in deliv.get("sub_deliverables", []):
                y += Inches(0.65)
                _add_textbox(slide, f"   • {sub}", Inches(0.8), y,
                             Inches(11.5), Inches(0.3), font_size=9.5, color=MID_GRAY)
            y += Inches(0.7)


def _add_approach_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Approach & Methodology"))

    _add_textbox(slide, data.get("approach_narrative", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.8),
                 font_size=11, color=DARK_GRAY)

    y = Inches(2.5)
    for i, step in enumerate(data.get("methodology_steps", []), 1):
        _add_textbox(slide, f"{i}. {step.get('step', '')}", Inches(0.5), y,
                     Inches(12.3), Inches(0.35), font_size=11, bold=True, color=RED)
        activities = "  ·  ".join(step.get("activities", []))
        _add_textbox(slide, activities, Inches(0.8), y + Inches(0.35),
                     Inches(12), Inches(0.3), font_size=9.5, color=DARK_GRAY)
        y += Inches(0.85)


def _add_governance_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Engagement Governance"))

    _add_textbox(slide, data.get("governance_structure", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.8),
                 font_size=11, color=DARK_GRAY)

    _add_bullet_column(slide, "Reporting Cadence", data.get("reporting_cadence", []),
                       Inches(0.5), Inches(2.5), Inches(5.5))

    _add_textbox(slide, "Quality Assurance", Inches(7), Inches(2.4),
                 Inches(5.8), Inches(0.35), font_size=11, bold=True, color=RED)
    _add_textbox(slide, data.get("quality_assurance", ""), Inches(7), Inches(2.8),
                 Inches(5.8), Inches(1.2), font_size=10.5, color=DARK_GRAY)


def _add_team_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Project Team"))

    _add_textbox(slide, data.get("team_narrative", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 font_size=11, color=DARK_GRAY)

    members = data.get("team_members", [])
    cols = min(len(members), 3)
    col_w = Inches(12.3 / max(cols, 1))
    for i, member in enumerate(members[:4]):
        x = Inches(0.5) + i * col_w
        _add_textbox(slide, member.get("role", ""), x, Inches(2.4),
                     col_w, Inches(0.35), font_size=11, bold=True, color=RED)
        _add_textbox(slide, member.get("relevant_experience", ""), x, Inches(2.8),
                     col_w, Inches(0.4), font_size=9.5, color=MID_GRAY)
        resp_text = "\n".join([f"• {r}" for r in member.get("responsibilities", [])])
        _add_textbox(slide, resp_text, x, Inches(3.3), col_w, Inches(2.0),
                     font_size=9.5, color=DARK_GRAY)


def _add_timeline_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Engagement Timeline"))

    _add_textbox(slide, f"Total Duration: {data.get('total_duration', '')}",
                 Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
                 font_size=11, bold=True, color=RED)

    y = Inches(2.1)
    for phase in data.get("phases", []):
        _add_textbox(slide, f"{phase.get('phase', '')}  —  {phase.get('duration', '')}",
                     Inches(0.5), y, Inches(12.3), Inches(0.35),
                     font_size=11, bold=True, color=DARK_GRAY)
        milestones = "  ▸  ".join(phase.get("key_milestones", []))
        _add_textbox(slide, milestones, Inches(0.8), y + Inches(0.35),
                     Inches(12), Inches(0.3), font_size=9.5, color=MID_GRAY)
        y += Inches(0.9)


def _add_experience_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Our Relevant Experience"))

    _add_textbox(slide, data.get("narrative", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 font_size=11, color=DARK_GRAY)

    case_studies = data.get("case_studies", [])
    cols = min(len(case_studies), 3)
    col_w = Inches(12.3 / max(cols, 1))
    for i, cs in enumerate(case_studies[:3]):
        x = Inches(0.5) + i * col_w
        # Box background
        box = slide.shapes.add_shape(1, x, Inches(2.4), col_w - Inches(0.2), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()
        _add_textbox(slide, cs.get("client_type", ""), x + Inches(0.1), Inches(2.5),
                     col_w - Inches(0.3), Inches(0.35), font_size=9, bold=True, color=MID_GRAY)
        _add_textbox(slide, cs.get("engagement", ""), x + Inches(0.1), Inches(2.9),
                     col_w - Inches(0.3), Inches(0.6), font_size=11, bold=True, color=RED)
        _add_textbox(slide, cs.get("our_role", ""), x + Inches(0.1), Inches(3.6),
                     col_w - Inches(0.3), Inches(1.0), font_size=9.5, color=DARK_GRAY)
        _add_textbox(slide, f"Outcome: {cs.get('outcome', '')}", x + Inches(0.1), Inches(4.7),
                     col_w - Inches(0.3), Inches(0.8), font_size=9.5, bold=True, color=RED)


def _add_why_protiviti(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Why Protiviti"))

    _add_textbox(slide, data.get("headline", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 font_size=14, bold=True, color=RED)

    y = Inches(2.3)
    for reason in data.get("reasons", []):
        _add_textbox(slide, f"◆  {reason.get('reason', '')}", Inches(0.5), y,
                     Inches(12.3), Inches(0.35), font_size=11, bold=True, color=DARK_GRAY)
        _add_textbox(slide, reason.get("detail", ""), Inches(0.9), y + Inches(0.35),
                     Inches(12), Inches(0.4), font_size=10, color=MID_GRAY)
        y += Inches(0.95)


def _add_assumptions_slide(prs, data: dict):
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, data.get("title", "Key Assumptions & Dependencies"))

    _add_bullet_column(slide, "Key Assumptions", data.get("assumptions", []),
                       Inches(0.5), Inches(1.6), Inches(4))
    _add_bullet_column(slide, "Client Dependencies", data.get("client_dependencies", []),
                       Inches(4.8), Inches(1.6), Inches(4))
    _add_bullet_column(slide, "Out of Scope", data.get("out_of_scope", []),
                       Inches(9.0), Inches(1.6), Inches(4))


def _add_closing_slide(prs, client_name: str):
    slide = _blank_slide(prs)
    _set_bg(slide, RED)
    _add_textbox(slide, "Thank You", Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, f"We look forward to partnering with {client_name}",
                 Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "PROTIVITI MIDDLE EAST | REAL ESTATE & INFRASTRUCTURE PRACTICE",
                 Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 font_size=10, color=RGBColor(0xFF, 0x99, 0x99), align=PP_ALIGN.CENTER)


def _add_slide_header(slide, title: str):
    """Add a standard slide title with red underline."""
    _add_textbox(slide, title, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 font_size=20, bold=True, color=RED)
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def _add_bullet_column(slide, heading: str, bullets: list, x, y, width):
    _add_textbox(slide, heading, x, y, width, Inches(0.35),
                 font_size=11, bold=True, color=RED)
    text = "\n".join([f"• {b}" for b in bullets])
    _add_textbox(slide, text, x, y + Inches(0.4), width, Inches(4.5),
                 font_size=10, color=DARK_GRAY)


def _add_client_overview_slide(prs, client_profile: dict, rfp: dict):
    """Client overview slide powered by web research."""
    if not client_profile or client_profile.get("error"):
        return
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, f"About {rfp.get('client_name', 'the Client')}")

    # Overview narrative
    _add_textbox(slide,
                 client_profile.get("organization_overview", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.9),
                 font_size=11, color=DARK_GRAY)

    # Three columns: Mandate | Recent Initiatives | Leadership Priorities
    _add_textbox(slide, "Core Mandate", Inches(0.5), Inches(2.6),
                 Inches(3.8), Inches(0.35), font_size=10, bold=True, color=RED)
    _add_textbox(slide, client_profile.get("mandate_and_role", ""),
                 Inches(0.5), Inches(3.0), Inches(3.8), Inches(1.5),
                 font_size=9.5, color=DARK_GRAY)

    initiatives = client_profile.get("recent_initiatives", [])
    init_text = "\n".join([f"• {i.get('initiative', i) if isinstance(i, dict) else i}"
                           for i in initiatives[:4]])
    _add_textbox(slide, "Recent Strategic Initiatives", Inches(4.6), Inches(2.6),
                 Inches(4), Inches(0.35), font_size=10, bold=True, color=RED)
    _add_textbox(slide, init_text, Inches(4.6), Inches(3.0),
                 Inches(4), Inches(2.5), font_size=9.5, color=DARK_GRAY)

    _add_bullet_column(slide, "Leadership Priorities",
                       client_profile.get("leadership_priorities", []),
                       Inches(9.0), Inches(2.6), Inches(4.0))

    # Key stats bar at bottom
    stats = client_profile.get("key_statistics", [])
    if stats:
        _add_textbox(slide, "  ·  ".join(stats[:3]),
                     Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
                     font_size=8.5, color=MID_GRAY)

    # Research attribution
    sources = client_profile.get("sources_used", [])
    if sources:
        _add_textbox(slide, f"Sources: {', '.join(sources[:3])}",
                     Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
                     font_size=7, color=RGBColor(0xBB, 0xBB, 0xBB))


def _add_sector_context_slide(prs, sector_context: dict, rfp: dict):
    """Sector landscape slide powered by web research."""
    if not sector_context or sector_context.get("error"):
        return
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide,
                      f"{rfp.get('sector', 'Sector')} Landscape — {rfp.get('geography', 'UAE')}")

    _add_textbox(slide, sector_context.get("sector_overview", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.7),
                 font_size=11, color=DARK_GRAY)

    # Transformation drivers (left)
    _add_bullet_column(slide, "Transformation Drivers",
                       sector_context.get("transformation_drivers", []),
                       Inches(0.5), Inches(2.4), Inches(4))

    # National programs (centre)
    programs = sector_context.get("national_programs", [])
    prog_text = "\n".join([
        f"• {p.get('program', p) if isinstance(p, dict) else p}"
        for p in programs[:4]
    ])
    _add_textbox(slide, "National Programs & Initiatives",
                 Inches(4.8), Inches(2.4), Inches(4), Inches(0.35),
                 font_size=11, bold=True, color=RED)
    _add_textbox(slide, prog_text, Inches(4.8), Inches(2.8),
                 Inches(4), Inches(2.5), font_size=9.5, color=DARK_GRAY)

    # Sector challenges (right)
    _add_bullet_column(slide, "Key Sector Challenges",
                       sector_context.get("sector_challenges", []),
                       Inches(9.0), Inches(2.4), Inches(4.0))

    # Stats
    stats = sector_context.get("market_statistics", [])
    if stats:
        _add_textbox(slide, "  ·  ".join(stats[:2]),
                     Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                     font_size=8.5, color=MID_GRAY)

    sources = sector_context.get("sources_used", [])
    if sources:
        _add_textbox(slide, f"Sources: {', '.join(sources[:3])}",
                     Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
                     font_size=7, color=RGBColor(0xBB, 0xBB, 0xBB))


def _add_value_add_slide(prs, value_add: dict):
    """Leading practice recommendations slide."""
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_slide_header(slide, value_add.get("title", "Our Recommended Enhancements"))

    _add_textbox(slide, value_add.get("intro", ""),
                 Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 font_size=10.5, color=DARK_GRAY)

    # Included in scope (green accent)
    included = value_add.get("included_in_scope", [])
    if included:
        _add_textbox(slide, "Included in Scope — Leading Practice Additions",
                     Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.35),
                     font_size=10, bold=True, color=RGBColor(0x1A, 0x7A, 0x3C))
        y = Inches(2.75)
        for item in included[:3]:
            _add_textbox(slide, f"✓  {item.get('title', '')}",
                         Inches(0.5), y, Inches(5.8), Inches(0.3),
                         font_size=10, bold=True, color=DARK_GRAY)
            _add_textbox(slide, item.get("talking_point", ""),
                         Inches(0.8), y + Inches(0.3), Inches(5.5), Inches(0.35),
                         font_size=9, color=MID_GRAY)
            y += Inches(0.8)

    # Optional add-ons (right column)
    optional = value_add.get("optional_addons", [])
    if optional:
        _add_textbox(slide, "Optional Enhancements",
                     Inches(7.0), Inches(2.3), Inches(5.8), Inches(0.35),
                     font_size=10, bold=True, color=RED)
        y = Inches(2.75)
        for item in optional[:3]:
            fee = item.get("fee_usd", 0)
            _add_textbox(slide,
                         f"+ {item.get('title', '')}  |  USD {fee:,.0f}",
                         Inches(7.0), y, Inches(5.8), Inches(0.3),
                         font_size=10, bold=True, color=DARK_GRAY)
            _add_textbox(slide, item.get("benefit", ""),
                         Inches(7.2), y + Inches(0.3), Inches(5.5), Inches(0.35),
                         font_size=9, color=MID_GRAY)
            y += Inches(0.8)

    # Future phases
    future = value_add.get("future_phases", [])
    if future:
        future_text = "  ·  ".join([f.get("title", "") for f in future[:3]])
        _add_textbox(slide, f"Natural Follow-On:  {future_text}",
                     Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
                     font_size=9.5, color=MID_GRAY)
