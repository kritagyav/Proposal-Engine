"""
Reads past proposals (PPT/PDF/DOCX) from the proposals folder
and extracts text content for indexing.
"""
import os
from pathlib import Path
from pptx import Presentation
import pdfplumber
from docx import Document


def extract_text_from_pptx(file_path: str) -> str:
    """Extract all text from a PowerPoint file, slide by slide."""
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract all text from a PDF file."""
    text_pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                text_pages.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(text_pages)


def extract_text_from_docx(file_path: str) -> str:
    """Extract all text from a Word document."""
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text(file_path: str) -> str | None:
    """Route to correct extractor based on file extension."""
    ext = Path(file_path).suffix.lower()
    try:
        if ext in (".pptx", ".ppt"):
            return extract_text_from_pptx(file_path)
        elif ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return extract_text_from_docx(file_path)
    except Exception as e:
        print(f"  Error reading {file_path}: {e}")
    return None


def load_all_proposals(proposals_folder: str) -> list[dict]:
    """
    Scan proposals folder and extract text from all supported files.
    Returns list of dicts: {filename, text, file_path}
    """
    supported = {".pptx", ".ppt", ".pdf", ".docx", ".doc"}
    proposals = []

    folder = Path(proposals_folder)
    if not folder.exists():
        print(f"Proposals folder not found: {proposals_folder}")
        return proposals

    files = [f for f in folder.iterdir() if f.suffix.lower() in supported]
    print(f"Found {len(files)} proposal files in {proposals_folder}")

    for file in files:
        print(f"  Reading: {file.name}")
        text = extract_text(str(file))
        if text and len(text.strip()) > 100:
            proposals.append({
                "filename": file.name,
                "file_path": str(file),
                "text": text,
            })
        else:
            print(f"  Skipped (no content): {file.name}")

    print(f"\nSuccessfully extracted {len(proposals)} proposals.")
    return proposals
